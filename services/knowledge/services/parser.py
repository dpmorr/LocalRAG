"""
Document parsing service
Supports: PDF, DOCX, TXT, MD, HTML, CSV, XLSX, XLS, RTF, ODT, PPTX, Images (OCR)
"""
import logging
from pathlib import Path
from typing import Dict, Any
from pypdf import PdfReader
from docx import Document
import io
from bs4 import BeautifulSoup
import re
import json
import csv

logger = logging.getLogger(__name__)


class DocumentParser:
    """Parse various document formats to text"""

    @staticmethod
    async def parse(file_content: bytes, content_type: str, filename: str) -> Dict[str, Any]:
        """
        Parse document to text - supports many formats

        Returns:
            {
                "text": str,
                "metadata": dict
            }
        """
        try:
            # Get file extension
            ext = filename.lower().split('.')[-1] if '.' in filename else ''

            # PDF
            if content_type == "application/pdf" or ext == 'pdf':
                return await DocumentParser._parse_pdf(file_content, filename)

            # Microsoft Word
            elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or ext == 'docx':
                return await DocumentParser._parse_docx(file_content, filename)

            # Plain text and markdown
            elif content_type in ["text/plain", "text/markdown"] or ext in ['txt', 'md']:
                return await DocumentParser._parse_text(file_content, filename)

            # CSV
            elif content_type == "text/csv" or ext == 'csv':
                return await DocumentParser._parse_csv(file_content, filename)

            # Excel (XLSX, XLS)
            elif content_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"] or ext in ['xlsx', 'xls']:
                return await DocumentParser._parse_excel(file_content, filename)

            # JSON
            elif content_type == "application/json" or ext == 'json':
                return await DocumentParser._parse_json(file_content, filename)

            # HTML
            elif content_type == "text/html" or ext in ['html', 'htm']:
                return await DocumentParser._parse_html_file(file_content, filename)

            # PowerPoint
            elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or ext == 'pptx':
                return await DocumentParser._parse_pptx(file_content, filename)

            # RTF
            elif content_type == "application/rtf" or ext == 'rtf':
                return await DocumentParser._parse_rtf(file_content, filename)

            # Images (OCR)
            elif content_type.startswith("image/") or ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                return await DocumentParser._parse_image(file_content, filename, content_type)

            else:
                # Try to decode as text
                try:
                    text = file_content.decode('utf-8')
                    return {
                        "text": text,
                        "metadata": {"source": filename, "format": "text"}
                    }
                except:
                    raise ValueError(f"Unsupported content type: {content_type} (file: {filename})")

        except Exception as e:
            logger.error(f"Failed to parse document {filename}: {e}")
            raise

    @staticmethod
    async def _parse_pdf(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse PDF file"""
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)

        text_parts = []
        for page_num, page in enumerate(reader.pages):
            text = page.extract_text()
            text_parts.append(text)

        full_text = "\n\n".join(text_parts)

        metadata = {
            "pages": len(reader.pages),
            "source": filename,
            "format": "pdf"
        }

        # Try to extract PDF metadata
        if reader.metadata:
            metadata.update({
                "title": reader.metadata.get("/Title", ""),
                "author": reader.metadata.get("/Author", ""),
                "subject": reader.metadata.get("/Subject", ""),
            })

        return {
            "text": full_text,
            "metadata": metadata
        }

    @staticmethod
    async def _parse_docx(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse DOCX file"""
        docx_file = io.BytesIO(content)
        doc = Document(docx_file)

        # Extract all paragraphs
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
        full_text = "\n\n".join(text_parts)

        metadata = {
            "paragraphs": len(doc.paragraphs),
            "source": filename,
            "format": "docx"
        }

        # Extract core properties if available
        if doc.core_properties:
            metadata.update({
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "subject": doc.core_properties.subject or "",
            })

        return {
            "text": full_text,
            "metadata": metadata
        }

    @staticmethod
    async def _parse_text(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse plain text or markdown file"""
        text = content.decode("utf-8")

        metadata = {
            "source": filename,
            "format": "markdown" if filename.endswith(".md") else "text"
        }

        return {
            "text": text,
            "metadata": metadata
        }

    @staticmethod
    async def parse_html(html_content: str, url: str) -> Dict[str, Any]:
        """
        Parse HTML content from URL to clean markdown-like text

        Args:
            html_content: Raw HTML content
            url: Source URL

        Returns:
            {
                "text": str,
                "metadata": dict
            }
        """
        try:
            soup = BeautifulSoup(html_content, 'html.parser')

            # Remove script and style elements
            for script in soup(["script", "style", "nav", "header", "footer"]):
                script.decompose()

            # Extract title
            title = ""
            if soup.title:
                title = soup.title.string.strip()
            elif soup.find('h1'):
                title = soup.find('h1').get_text().strip()

            # Extract main content
            # Try to find main content area
            main_content = None
            for tag in ['article', 'main', '[role="main"]']:
                main_content = soup.find(tag)
                if main_content:
                    break

            if not main_content:
                main_content = soup.find('body') or soup

            # Extract text with structure
            text_parts = []

            for element in main_content.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'blockquote']):
                text = element.get_text().strip()
                if not text:
                    continue

                # Add markdown-style formatting
                if element.name.startswith('h'):
                    level = int(element.name[1])
                    text = f"{'#' * level} {text}"
                elif element.name == 'li':
                    text = f"- {text}"
                elif element.name == 'blockquote':
                    text = f"> {text}"

                text_parts.append(text)

            full_text = "\n\n".join(text_parts)

            # Extract metadata
            metadata = {
                "source": url,
                "format": "article",
                "title": title
            }

            # Try to extract author and date
            author_meta = soup.find('meta', attrs={'name': 'author'}) or \
                         soup.find('meta', attrs={'property': 'article:author'})
            if author_meta:
                metadata["author"] = author_meta.get('content', '')

            date_meta = soup.find('meta', attrs={'property': 'article:published_time'}) or \
                       soup.find('meta', attrs={'name': 'date'})
            if date_meta:
                metadata["published_date"] = date_meta.get('content', '')

            # Extract description
            desc_meta = soup.find('meta', attrs={'name': 'description'}) or \
                       soup.find('meta', attrs={'property': 'og:description'})
            if desc_meta:
                metadata["description"] = desc_meta.get('content', '')

            return {
                "text": full_text,
                "metadata": metadata
            }

        except Exception as e:
            logger.error(f"Failed to parse HTML from {url}: {e}")
            raise

    @staticmethod
    async def _parse_csv(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse CSV file"""
        text_io = io.StringIO(content.decode('utf-8'))
        reader = csv.DictReader(text_io)
        
        rows = list(reader)
        text_parts = []
        
        # Convert CSV to readable format
        for i, row in enumerate(rows):
            row_text = f"Row {i+1}:\n"
            for key, value in row.items():
                row_text += f"  {key}: {value}\n"
            text_parts.append(row_text)
        
        return {
            "text": "\n".join(text_parts),
            "metadata": {
                "source": filename,
                "format": "csv",
                "rows": len(rows),
                "columns": len(rows[0].keys()) if rows else 0
            }
        }

    @staticmethod
    async def _parse_excel(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse Excel file (XLSX/XLS)"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content))
            
            text_parts = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            
            return {
                "text": "\n".join(text_parts),
                "metadata": {
                    "source": filename,
                    "format": "excel",
                    "sheets": len(wb.sheetnames)
                }
            }
        except ImportError:
            return {
                "text": f"Excel file: {filename} (Excel parser not available)",
                "metadata": {"source": filename, "format": "excel", "error": "parser_unavailable"}
            }

    @staticmethod
    async def _parse_json(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse JSON file"""
        data = json.loads(content.decode('utf-8'))
        
        # Convert JSON to readable text
        text = json.dumps(data, indent=2)
        
        return {
            "text": text,
            "metadata": {
                "source": filename,
                "format": "json"
            }
        }

    @staticmethod
    async def _parse_html_file(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse HTML file"""
        html_content = content.decode('utf-8')
        return await DocumentParser.parse_html(html_content, filename)

    @staticmethod
    async def _parse_pptx(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse PowerPoint file"""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            
            text_parts = []
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"Slide {i+1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            return {
                "text": "\n\n".join(text_parts),
                "metadata": {
                    "source": filename,
                    "format": "pptx",
                    "slides": len(prs.slides)
                }
            }
        except ImportError:
            return {
                "text": f"PowerPoint file: {filename} (PowerPoint parser not available)",
                "metadata": {"source": filename, "format": "pptx", "error": "parser_unavailable"}
            }

    @staticmethod
    async def _parse_rtf(content: bytes, filename: str) -> Dict[str, Any]:
        """Parse RTF file"""
        try:
            from striprtf.striprtf import rtf_to_text
            text = rtf_to_text(content.decode('utf-8', errors='ignore'))
            
            return {
                "text": text,
                "metadata": {
                    "source": filename,
                    "format": "rtf"
                }
            }
        except ImportError:
            # Fallback: try to extract text naively
            text = content.decode('utf-8', errors='ignore')
            # Remove RTF control words
            text = re.sub(r'\[a-z]+\d*\s?', '', text)
            text = re.sub(r'[{}]', '', text)
            
            return {
                "text": text,
                "metadata": {
                    "source": filename,
                    "format": "rtf",
                    "note": "Basic RTF parsing"
                }
            }

    @staticmethod
    async def _parse_image(content: bytes, filename: str, content_type: str) -> Dict[str, Any]:
        """Parse image file using OCR"""
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            
            return {
                "text": text,
                "metadata": {
                    "source": filename,
                    "format": "image",
                    "type": content_type,
                    "ocr": "tesseract"
                }
            }
        except ImportError:
            return {
                "text": f"Image file: {filename} (OCR not available - install pytesseract and PIL)",
                "metadata": {
                    "source": filename,
                    "format": "image",
                    "type": content_type,
                    "error": "ocr_unavailable"
                }
            }
